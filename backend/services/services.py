from genericpath import exists
from model.mongo import mongoModel
from engine.youtube_crawler import youtube
from engine.google_crawler import google
import os
import pdfplumber
from pptx import Presentation
from engine.keyword_extraction import extract_keyword

mm = mongoModel()
y = youtube()
g = google()

guest_links = []

def verify(uid, pwd):
    userinfo = mm.get_user(uid)
    if userinfo is not None and userinfo["pwd"] == pwd:
        return True
    else:
        return False

def signup(uid, pwd):
    if check_redundancy(uid):
        user_info = {
        "uid" : uid,
        "pwd" : pwd,
        "files" : [{"type":"DIR","parent":-1}],
        }
        mm.set_user(user_info)
        return True
    else:
        return False
    
def check_redundancy(uid):
    user = mm.get_user(uid)
    if user is None:
        return True
    else:
        return False

def files(uid, id):
    return mm.get_files(uid, id)

def file(uid, id):
    return mm.get_file(uid, id)

def check_dir(uid, id):
    if id == 0:
        return files(uid, id)
    else:
        file_info = file(uid, id)
        if file_info["type"] == "DIR":
            return files(uid, id)
        else:
            return file_info

def upload(uid, filename, cdi):
    file_path = "./static/files/"+filename
    file_info = {
        "name" : filename.split(".")[0],
        "type" : filename.split(".")[-1].upper(),
        "parent" : int(cdi),
    }
    with open(file_path, "rb") as f:
        file_info = mm.put_file(uid, f, file_info)
        if file_info:
            return file_info
        else:
            return False
    

def makedir(uid, dirname, cdi):
    dir_info = {
        "name" : dirname,
        "type" : "DIR",
        "parent" : int(cdi)
    }
    if mm.put_file(uid, None, dir_info):
        return True
    else:
        return False

def read_file(file_info):
    ft = file_info["type"].lower()
    if "path" in file_info:
        path = file_info["path"]
    else:
        path = download(file_info)
        
    if ft == "pdf":
        pr = pdfplumber.open(path)
    elif ft == "ppt" or ft == "pptx":
        pr = Presentation(path)
    
    return pr, ft
    
def read_page(pr, pagenum, ft):
    if ft == "pdf":
        text = pr.pages[pagenum].extract_text()
        keywords = extract_keyword(text)
        links = y.get_youtube_links(keywords)
        return links
    elif ft == "ppt" or ft == "pptx":
        text_runs = []
        slides = pr.slides
        for shape in slides[pagenum].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text_runs.append(run.text)
        return text_runs

def download(file_info):
    path = "./static/files/"+file_info["name"]+"."+file_info["type"].lower()
    if 'fileid' in file_info:
        file = mm.get_file_from_fs(file_info['fileid'])
        with open(path, "wb") as f:
            f.write(file)
        return path
    else:
        return False

def process_file(uid, file_info):
    pr, ft = read_file(file_info)
    pagelen = len(pr.pages)
    if uid is not None:   
        data = []
        for i in range(pagelen):
            data.append({"links" : read_page(pr, i, ft)})
        if mm.set_links(uid, file_info["id"], data):
            return True
        else:
            return False
    #Guest
    else:
        guest_links.clear()
        for i in range(pagelen):
            guest_links.append({"links" : read_page(pr, i, ft)})
        
def get_link(uid, id, pagenum):
    if uid is not None:
        link = mm.get_file(uid, id)["data"][int(pagenum)]
        if link == None:
            return False
        else:
            return link
    #Guest
    else:
        return guest_links[int(pagenum)]